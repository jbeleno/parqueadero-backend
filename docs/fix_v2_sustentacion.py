"""
Fix-v2 con índices CORRECTOS (post-eliminación) y estructura por slide real.
"""
import pathlib
from copy import deepcopy
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

ROOT = pathlib.Path("/Users/jesus/Desktop/parqueaderos-api")
OUT = ROOT / "parqueaderos-docs-referencia/SUBIR_A_DRIVE/Sustentacion_Jesus.pptx"


def set_paragraph_text(paragraph, new_text):
    if not paragraph.runs:
        run_xml = '<a:r xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:rPr lang="es-CO" dirty="0"/><a:t></a:t></a:r>'
        new_r = etree.fromstring(run_xml)
        paragraph._p.append(new_r)
    paragraph.runs[0].text = new_text
    for r in paragraph.runs[1:]:
        r._r.getparent().remove(r._r)


def set_shape_lines(shape, lines):
    """Setea línea-por-párrafo (modo simple)."""
    if not shape.has_text_frame:
        return
    paras = list(shape.text_frame.paragraphs)
    for i, line in enumerate(lines):
        if i < len(paras):
            set_paragraph_text(paras[i], line)
    for i in range(len(lines), len(paras)):
        set_paragraph_text(paras[i], "")


def replace_text_frame_multi(shape, lines):
    """Reemplaza completamente el text_frame con N párrafos clonando el formato del primero."""
    if not shape.has_text_frame:
        return
    paras = list(shape.text_frame.paragraphs)
    if not paras:
        return
    template_p = paras[0]._p
    first_run = template_p.find(qn("a:r"))
    if first_run is not None and first_run.find(qn("a:rPr")) is not None:
        rPr_xml = etree.tostring(first_run.find(qn("a:rPr")))
    else:
        rPr_xml = b'<a:rPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" lang="es-CO" dirty="0"/>'
    pPr_el = template_p.find(qn("a:pPr"))
    pPr_xml = etree.tostring(pPr_el) if pPr_el is not None else None
    txBody = template_p.getparent()
    for p in list(txBody.findall(qn("a:p"))):
        txBody.remove(p)
    for line in lines:
        new_p = etree.SubElement(txBody, qn("a:p"))
        if pPr_xml is not None:
            new_p.insert(0, etree.fromstring(pPr_xml))
        new_r = etree.SubElement(new_p, qn("a:r"))
        new_r.append(etree.fromstring(rPr_xml))
        new_t = etree.SubElement(new_r, qn("a:t"))
        new_t.text = line


def remove_picture_by_size(slide, min_w, max_w):
    n = 0
    for shape in list(slide.shapes):
        if shape.shape_type != 13:
            continue
        if min_w <= shape.width <= max_w:
            shape._element.getparent().remove(shape._element)
            n += 1
    return n


# ─────────────────────────────────────────────────────────────────────────────
# Contenido por slide (índices POST-eliminación)
# ─────────────────────────────────────────────────────────────────────────────

# Slides que usan shape-por-bullet (estructura del template Linda)
PER_SHAPE = {
    2: {  # Título
        0: ["Backend transaccional con visión artificial",
            "para gestión integral de parqueaderos"],
        1: ["Reconocimiento automático de placas (ALPR) — YOLOv11n + PaddleOCR"],
        2: ["JESÚS BELEÑO", "PROGRAMA DE INGENIERÍA DE SOFTWARE"],
    },
    4: {  # Introducción (5 shapes: título + 4 bullets)
        0: ["Introducción"],
        1: ["Backend transaccional que actúa como columna vertebral del sistema de gestión de parqueaderos."],
        2: ["Spring Boot 3 + Java 21 con ~245 endpoints REST, JWT, WebSocket STOMP y PostgreSQL 16 + PostGIS."],
        3: ["Sidecar Python (FastAPI) con YOLOv11n y PaddleOCR para reconocimiento automático de placas."],
        4: ["Desplegado como contenedores Docker en Dokploy con HTTPS por Traefik."],
    },
    5: {  # Problema (6 shapes)
        0: ["Planteamiento del Problema"],
        1: ["Un parqueadero requiere registrar entradas/salidas, cobrar tarifas, coordinar barreras y cámaras, y prevenir fraudes."],
        2: ["Los sistemas comerciales son cerrados y no integran visión por computador para placas colombianas."],
        3: ["Tampoco soportan multi-tenant con resoluciones DIAN, IVA desagregado o Modelo B regulatorio."],
        4: ["Reto técnico: acoplar deep learning Python con un backend transaccional Java sin comprometer la estabilidad."],
        5: ["Solución: backend Java + sidecar Python comunicados por HTTP, con tolerancia a fallos del sidecar."],
    },
    6: {  # Objetivos (10 shapes: título + bullets + labels)
        0: ["Objetivos"],
        4: ["Objetivo General"],
        1: ["Diseñar, construir y desplegar el backend transaccional y el módulo de visión artificial del sistema integral de gestión de parqueaderos."],
        5: ["Objetivos Específicos"],
        2: ["Diseñar el modelo de datos y la arquitectura multi-tenant con JWT, RBAC granular y auditoría universal."],
        6: ["Implementar tarifas Modelo B (IVA, franjas, tope) y flujos de tickets manuales, OCR y reservas."],
        8: ["Entrenar y desplegar el módulo OCR (YOLOv11n + PaddleOCR) como sidecar tolerante a fallos integrado por WebSocket."],
    },
    7: {  # Justificación (7 shapes: título + 3 pares label/bullet)
        0: ["Justificación"],
        1: ["La integración de visión por computador y eventos en tiempo real reduce costos operativos y aumenta la trazabilidad."],
        2: ["Arquitectura desacoplada"],
        3: ["Combina un sistema empresarial Java con el ecosistema de visión por computador de Python sin sacrificar estabilidad."],
        4: ["Tolerancia a fallos"],
        5: ["Degradación silenciosa: un fallo del sidecar OCR no interrumpe el flujo operativo."],
        6: ["Cumplimiento fiscal"],
        7: ["Modelo B colombiano con IVA desagregado y resoluciones DIAN listas para facturación electrónica."],
    },
    8: {  # Metodología
        0: ["Metodología"],
        1: ["Cinco fases del ciclo de vida del software, adaptadas a un proyecto unipersonal con entregas incrementales."],
        2: ["Análisis → Diseño → Implementación → Pruebas → Despliegue, con retroalimentación continua entre fases."],
    },
    9: {  # Análisis
        0: ["Fase de Análisis"],
        1: ["Identificación de actores: conductor, vigilante, ADMIN, SUPER_ADMIN, cámara IoT, sidecar OCR y SMTP."],
        2: ["41 requerimientos funcionales (RF-01 a RF-41) y 17 no funcionales con plantilla IEEE 29148."],
        3: ["41 historias de usuario en formato 'Como X, quiero Y, para Z' mapeadas 1:1 con los RF."],
        4: ["42 diagramas de casos de uso UML: uno general y 41 individuales, uno por cada RF."],
    },
    12: {  # Diseño (5 shapes)
        0: ["Fase de Diseño"],
        1: ["Arquitectura: cliente ↔ backend Spring Boot ↔ PostgreSQL + PostGIS; sidecar Python OCR por REST."],
        2: ["Modelo jerárquico Empresa → Parqueadero → Nivel → Sección → SubSección → PuntoParqueo, con PostGIS."],
        3: ["16 ERDs por dominio, 10 diagramas de estados y 20 de secuencia para los flujos críticos."],
        4: ["Multi-tenant scopado por empresa_id, soft-delete uniforme por archivado_en y auditoría AOP."],
    },
    16: {  # Implementación (10 shapes)
        0: ["Fase de Implementación"],
        1: ["Backend Spring Boot 3.5 + Java 21 en 20 paquetes: 34 controladores, 57 servicios, 68 repositorios, 68 entidades."],
        2: ["Sidecar Python: YOLOv11n (mAP50=0.987) recorta placa; PaddleOCR v2 con 3 preprocesos + voting (~2.5 s/img)."],
        3: ["API documentada en /swagger-ui.html; respuestas ApiResponse<T>; WebSocket STOMP en /topic y /queue."],
        4: ["Migraciones idempotentes en data.sql para reaplicar el seed sin romper datos preexistentes."],
        5: ["Auditoría append-only por aspecto AOP @Auditable, sin contaminar la lógica de negocio."],
    },
    18: {  # Pruebas (5 shapes)
        0: ["Fase de Pruebas"],
        1: ["139 pruebas unitarias en 24 archivos: TarifaCalculator (40), Ticket (26), Pago (13), Convenio (14), Backfill (7)."],
        2: ["Pruebas de integración con PostgreSQL real para listeners y jobs programados."],
        3: ["End-to-end desde Postman: registro → login → ticket → factura → pago → caja, validando WebSocket."],
        4: ["Benchmark OCR sobre 60 imágenes: PaddleOCR v2 ganó con 100% de consistencia frente a 7 alternativas."],
    },
    19: {  # Despliegue (12 shapes: 5 bullets + 4 cards de proceso)
        0: ["Fase de Despliegue"],
        1: ["Docker multi-stage: builder con Maven cachea dependencias; imagen final eclipse-temurin:21-jre-alpine (~250 MB)."],
        2: ["Despliegue automático en Dokploy: push a main → GitHub Actions construye → webhook → docker compose up."],
        3: ["HTTPS por Traefik con certificado auto-renovado; Swagger UI, 245 endpoints REST y WebSocket en el mismo dominio."],
        4: ["Volúmenes persistentes, red Docker privada y observabilidad con Spring Actuator."],
        # 5-8: cards Build/Distribución/Pruebas/Estabilización — mantenemos
    },
    20: {  # Alcance y Limitaciones (11 shapes)
        0: ["Alcance y Limitaciones"],
        1: ["Alcance"],
        2: ["~245 endpoints REST: autenticación, multi-tenant, tickets, reservas, facturación DIAN, pagos, caja, suscripciones, auditoría."],
        4: ["Sidecar OCR con YOLOv11n + PaddleOCR v2, latencia ~2.5 s/img y consistencia del 100% sobre el benchmark."],
        5: ["Limitaciones"],
        6: ["El alcance es backend + sidecar OCR; el cliente móvil y web son desarrollos paralelos no incluidos."],
        7: ["Tarifas Modelo B aditivo y de reemplazo; no contempla surge pricing ni descuentos por fidelidad."],
        8: ["Montos viejos en double precision; campos DIAN nuevos en numeric(14,2). Migración pendiente."],
        9: [""],   # vaciar texto viejo
        10: [""],  # vaciar texto viejo
    },
    21: {  # Resultados y Conclusiones (9 shapes)
        0: ["Resultados y Conclusiones"],
        1: ["Resultados"],
        2: ["71 tablas, 968 columnas, 78 FK, 379 CHECK y 30 UNIQUE. 139 @Test pasando. Sistema en producción."],
        4: ["mAP50=0.987 en YOLOv11n y consistencia del 100% en PaddleOCR v2 frente a 7 alternativas evaluadas."],
        5: ["Conclusiones"],
        6: ["Es viable acoplar deep learning a un sistema empresarial Java mediante HTTP delgado sin sacrificar estabilidad."],
        7: ["El Modelo B + IVA + tope + resoluciones DIAN dejan al sistema listo para certificación fiscal."],
        8: ["Auditoría AOP y RBAC granular soportan trazabilidad y aislamiento multi-tenant para operadores con múltiples sedes."],
    },
}

# Slides donde TODO el contenido va en shape[1] como múltiples párrafos
MULTI_PARAGRAPH = {
    10: [  # Ejemplo de RF — solo 2 shapes
        "Creación de ticket manual · RF-10 · Funcional · Prioridad Alta",
        "Descripción: crear ticket de entrada para un vehículo asignando un punto libre del parqueadero del operario.",
        "Entrada: POST /api/tickets con { parqueaderoId, puntoParqueoId, vehiculoId, tarifaId }",
        "Salida: HTTP 201 con TicketDTO en EN_CURSO + eventos WebSocket TICKET_CREADO y SPOT_STATUS_CHANGE.",
        "Precondición: punto LIBRE; operario OPERARIO_CAJA con caja ABIERTA.",
        "Postcondición: ticket EN_CURSO con bloqueo pesimista; evento publicado; operación auditada.",
        "Especificación completa de los 41 RF en el Anexo A.",
    ],
}


def main():
    prs = Presentation(OUT)

    print("═══ Aplicando contenido shape-por-bullet ═══")
    for slide_idx, shapes_map in PER_SHAPE.items():
        slide = prs.slides[slide_idx - 1]
        for shape_idx, lines in shapes_map.items():
            try:
                shape = slide.shapes[shape_idx]
                set_shape_lines(shape, lines)
                preview = lines[0][:50] if lines and lines[0] else "(vacío)"
                print(f"  slide {slide_idx:2d} · shape {shape_idx:2d} → {preview}")
            except IndexError:
                print(f"  ⚠ slide {slide_idx} no tiene shape {shape_idx}")

    print("\n═══ Aplicando contenido multi-párrafo (slide 10) ═══")
    for slide_idx, lines in MULTI_PARAGRAPH.items():
        slide = prs.slides[slide_idx - 1]
        # shape 1 es donde va el contenido
        replace_text_frame_multi(slide.shapes[1], lines)
        print(f"  slide {slide_idx} · shape 1 → {len(lines)} párrafos")

    print("\n═══ Eliminando logo E&T en slide 4 ═══")
    s4 = prs.slides[4 - 1]
    n = remove_picture_by_size(s4, min_w=500000, max_w=3500000)
    print(f"  slide 4: {n} imagen(es) eliminada(s)")

    prs.save(OUT)
    print(f"\n✓ Guardado: {OUT}")
    print(f"  Tamaño: {OUT.stat().st_size / 1024 / 1024:.1f} MB")


if __name__ == "__main__":
    main()
