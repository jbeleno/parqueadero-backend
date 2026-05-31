"""
Aplica la Opción B (layout por capas) al slide 14:
  1. Reemplaza image2.png con image2_optB.png
  2. Inserta la imagen en slide 14 manteniendo proporción y centrada
  3. Limpia las opciones descartadas (optA, optC1, optC2, optD)
"""
import pathlib
import shutil
from pptx import Presentation
from PIL import Image

ROOT = pathlib.Path("/Users/jesus/Desktop/parqueaderos-api")
OUT = ROOT / "parqueaderos-docs-referencia/SUBIR_A_DRIVE/presentacion/pptx/Sustentacion_Jesus.pptx"
FIG = ROOT / "docs/figures/main"

# Promover optB a image2.png oficial
shutil.copy(FIG / "image2_optB.png", FIG / "image2.png")
print(f"✓ image2_optB.png → image2.png")

# Limpiar opciones descartadas
for f in ["image2_optA.png", "image2_optC1.png", "image2_optC2.png", "image2_optD.png"]:
    p = FIG / f
    if p.exists():
        p.unlink()
        print(f"  🗑  {f} eliminado")

# Aplicar al slide 14
prs = Presentation(OUT)
sw, sh = prs.slide_width, prs.slide_height
s14 = prs.slides[14 - 1]

# Eliminar imagen actual
for shape in list(s14.shapes):
    if shape.shape_type == 13:
        shape._element.getparent().remove(shape._element)

# Re-insertar imagen nueva centrada
img_path = FIG / "image2.png"
with Image.open(img_path) as im:
    ratio = im.width / im.height

top_avail = int(sh * 0.18)
bottom_avail = int(sh * 0.96)
height_avail = bottom_avail - top_avail
max_width = int(sw * 0.95)

new_h = height_avail
new_w = int(new_h * ratio)
if new_w > max_width:
    new_w = max_width
    new_h = int(new_w / ratio)

left = (sw - new_w) // 2
top = top_avail + (height_avail - new_h) // 2

s14.shapes.add_picture(str(img_path), left, top, width=new_w, height=new_h)
print(f"  slide 14: L{left//100000} T{top//100000} W{new_w//100000} H{new_h//100000}")

prs.save(OUT)
print(f"\n✓ Guardado: {OUT}")
