"""
Fix-v3 del Sustentacion_Jesus.pptx.

Problemas detectados al renderizar v2 a PDF:
  - Slide 7: shapes 4↔5 y 7↔8 estaban INVERTIDOS (label vs contenido).
  - Slide 8: shape 3 quedó con texto viejo de Linda (duplicado).
  - Slide 9: shape 4 es un CARD DECORATIVO pequeño (W=28 H=6) en posición decorativa, no una bullet más. Texto largo desborda.
  - Slide 10: tiene una TABLA del template con la ficha de RF de Linda. Hay que LLENAR la tabla con datos reales y vaciar el shape 1 que metí mal.
  - Slide 12: shape 4 mismo problema que slide 9.
  - Slide 16: shapes 3, 6-9 son zonas decorativas / contenido viejo de Linda (móvil + web). Hay que limpiar.
  - Slide 17: imagen image5 está OK; shapes 1-5 vacíos OK.
  - Slide 18: shape 3 card decorativo.
  - Slide 19: shape 3 card decorativo; shapes 5-8 son los 4 pasos del despliegue (mantener).
  - Slide 20-21: estructura 2x2 con cards pequeños; mapeo distinto al que asumí.
"""
import pathlib
from copy import deepcopy
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

ROOT = pathlib.Path("/Users/jesus/Desktop/parqueaderos-api")
OUT = ROOT / "parqueaderos-docs-referencia/SUBIR_A_DRIVE/Sustentacion_Jesus.pptx"


def set_paragraph_text(paragraph, new_text):
    if not paragraph.runs:
        run_xml = '<a:r xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:rPr lang="es-CO" dirty="0"/><a:t></a:t></a:r>'
        new_r = etree.fromstring(run_xml)
        paragraph._p.append(new_r)
    paragraph.runs[0].text = new_text
    for r in paragraph.runs[1:]:
        r._r.getparent().remove(r._r)


def set_shape_lines(shape, lines):
    if not shape.has_text_frame:
        return
    paras = list(shape.text_frame.paragraphs)
    for i, line in enumerate(lines):
        if i < len(paras):
            set_paragraph_text(paras[i], line)
    for i in range(len(lines), len(paras)):
        set_paragraph_text(paras[i], "")


def clear_shape(shape):
    if shape.has_text_frame:
        set_shape_lines(shape, [""])


def set_table_cell(table, row_idx, col_idx, text):
    """Modifica el texto de una celda preservando formato."""
    cell = table.cell(row_idx, col_idx)
    tf = cell.text_frame
    paras = list(tf.paragraphs)
    if paras:
        set_paragraph_text(paras[0], text)
        # Eliminar párrafos adicionales
        first_p = paras[0]._p
        txBody = first_p.getparent()
        for p in list(txBody.findall(qn("a:p")))[1:]:
            txBody.remove(p)


def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


# ─────────────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation(OUT)

    # ── Slide 7: Justificación — invertir labels↔contenido ──
    print("═══ Fix slide 7: Justificación ═══")
    s7 = prs.slides[7 - 1]
    # shape 5 (T=26, W=28) es el LABEL del segundo card
    # shape 4 (T=33, W=115) es el CONTENIDO del segundo card
    set_shape_lines(s7.shapes[5], ["Tolerancia a fallos"])
    set_shape_lines(s7.shapes[4], ["Degradación silenciosa: un fallo del sidecar OCR no interrumpe el flujo operativo."])
    # shape 8 (T=43, W=28) es el LABEL del tercer card
    # shape 7 (T=50, W=115) es el CONTENIDO del tercer card
    set_shape_lines(s7.shapes[8], ["Cumplimiento fiscal"])
    set_shape_lines(s7.shapes[7], ["Modelo B colombiano con IVA desagregado y resoluciones DIAN listas para facturación electrónica."])
    print("  ✓ shapes 4↔5 y 7↔8 corregidos (label vs contenido)")

    # ── Slide 8: Metodología — vaciar shape 3 (duplicado) ──
    print("\n═══ Fix slide 8: Metodología ═══")
    s8 = prs.slides[8 - 1]
    clear_shape(s8.shapes[3])
    print("  ✓ shape 3 (texto viejo de Linda) vaciado")

    # ── Slide 9: Análisis — el shape 4 es un card decorativo pequeño ──
    print("\n═══ Fix slide 9: Análisis ═══")
    s9 = prs.slides[9 - 1]
    # El shape 4 (W=28 H=6) es un card en la esquina superior izquierda. Solo cabe texto MUY corto.
    set_shape_lines(s9.shapes[4], ["42 diagramas UML"])
    print("  ✓ shape 4 reducido a etiqueta corta '42 diagramas UML'")

    # ── Slide 10: Ejemplo RF — LLENAR la tabla y vaciar el bloque suelto ──
    print("\n═══ Fix slide 10: Ejemplo de RF ═══")
    s10 = prs.slides[10 - 1]
    # Vaciar el bloque de texto que metí mal (shape 1)
    clear_shape(s10.shapes[1])
    # Llenar la tabla (shape 2): 11 filas, 2 columnas
    tabla = s10.shapes[2].table
    # Cada fila tiene la etiqueta en col 0 (no la cambiamos) y el valor en col 1
    fila_valores = [
        "Creación de ticket manual",   # Título de Requerimiento
        "RF-10",                        # Código
        "[x] Funcional   [ ] No Funcional",   # Tipo
        "1.0",                          # Versión
        "Fase de análisis",             # Fuente
        "[x] Must Have (M)  [ ] Should Have  [ ] Could Have  [ ] Won't Have",  # Prioridad
        "[ ] Baja   [x] Media   [ ] Alta",   # Dificultad
        "RF-03 (Inicio de sesión con JWT)",  # Dependencia
        "Backend (Spring Boot) + Operario asignado al parqueadero",  # Responsables
        "POST /api/tickets con { parqueaderoId, puntoParqueoId, vehiculoId, tarifaId }",  # Entradas
        "[ ] Pendiente   [x] Implementado   [ ] Verificado   [ ] En Corrección   [ ] Aprobado",  # Estado
    ]
    for i, valor in enumerate(fila_valores):
        if i < len(tabla.rows):
            set_table_cell(tabla, i, 1, valor)
    print(f"  ✓ tabla llenada con {len(fila_valores)} filas del RF-10")

    # ── Slide 12: Diseño — vaciar card decorativo shape 4 ──
    print("\n═══ Fix slide 12: Diseño ═══")
    s12 = prs.slides[12 - 1]
    # El shape 4 (W=28 H=6) es card decorativo. Mi texto desbordó.
    set_shape_lines(s12.shapes[4], ["Diseño"])
    print("  ✓ shape 4 reducido a etiqueta 'Diseño'")

    # ── Slide 16: Implementación ──
    print("\n═══ Fix slide 16: Implementación ═══")
    s16 = prs.slides[16 - 1]
    # Estructura real:
    #   shape 3 (T=8, W=28) → card decorativo, texto corto
    #   shape 4 (T=31 central) → contenido columna central, texto corto
    #   shape 5 (T=30) → label/número de la columna central
    #   shape 6 (T=31 derecha) → label de la columna derecha "Sitio de Monitoreo Web" (Linda)
    #   shape 7 (T=30) → número columna derecha "2"
    #   shape 8 (T=40 izquierda) → lista de stack móvil (Linda)
    #   shape 9 (T=40 derecha) → lista de stack web (Linda)
    set_shape_lines(s16.shapes[3], ["Implementación"])
    set_shape_lines(s16.shapes[4], ["Backend Java"])
    set_shape_lines(s16.shapes[6], ["Sidecar Python OCR"])
    set_shape_lines(s16.shapes[8], [
        "Spring Boot 3.5 + Java 21",
        "PostgreSQL 16 + PostGIS",
        "Spring Security + JWT + RBAC",
        "WebSocket STOMP + Spring AOP",
    ])
    set_shape_lines(s16.shapes[9], [
        "Python 3.11 + FastAPI",
        "YOLOv11n (Ultralytics)",
        "PaddleOCR v2 + OpenCV",
        "Docker multi-stage",
    ])
    print("  ✓ slide 16 reorganizado en dos columnas: Backend Java | Sidecar Python")

    # ── Slide 18: Pruebas — vaciar card decorativo shape 3 ──
    print("\n═══ Fix slide 18: Pruebas ═══")
    s18 = prs.slides[18 - 1]
    set_shape_lines(s18.shapes[3], ["Pruebas"])
    print("  ✓ shape 3 → 'Pruebas'")

    # ── Slide 19: Despliegue — vaciar card decorativo shape 3 ──
    print("\n═══ Fix slide 19: Despliegue ═══")
    s19 = prs.slides[19 - 1]
    set_shape_lines(s19.shapes[3], ["Despliegue"])
    print("  ✓ shape 3 → 'Despliegue'")

    # ── Slide 20: Alcance y Limitaciones — reorganización completa ──
    print("\n═══ Fix slide 20: Alcance y Limitaciones ═══")
    s20 = prs.slides[20 - 1]
    # Estructura real (analizando coordenadas):
    #   shape 3 (T=13, W=8) → número "1" izquierdo
    #   shape 5 (T=13, W=8) → label "Limitaciones" + número
    #   shape 2 (T=14, W=13) → contenido pt1 izq (CHICO, texto corto)
    #   shape 4 (T=14, W=26) → contenido pt1 der
    #   shape 1 (T=23, W=56) → label "Alcance" (donde mi texto largo encaja MAL)
    #   shape 6 (T=34, W=56) → contenido pt2 izq (donde sí cabe)
    #   shape 7 (T=44, W=56) → contenido pt3 izq
    #   shape 8 (T=23, W=56) → contenido pt2 der
    # Estrategia: cuadros chicos = texto MUY corto; cuadros grandes = bullets.
    set_shape_lines(s20.shapes[3], ["1"])
    set_shape_lines(s20.shapes[5], ["2"])
    set_shape_lines(s20.shapes[2], ["245 endpoints REST"])
    set_shape_lines(s20.shapes[4], ["Sidecar OCR funcional"])
    set_shape_lines(s20.shapes[1], ["Alcance"])
    # En shape 1 metemos solo "Alcance" como label
    # Detalle del alcance va a shape 6
    set_shape_lines(s20.shapes[6], ["Backend con multi-tenant, tickets, reservas, facturación DIAN, suscripciones, convenios, caja, auditoría y RBAC granular."])
    set_shape_lines(s20.shapes[7], ["YOLOv11n con mAP50=0.987 y PaddleOCR v2 con 100% de consistencia sobre 60 imágenes de prueba."])
    set_shape_lines(s20.shapes[8], ["Limitaciones"])
    # shapes 9 y 10 son zonas extra (Linda tenía más cuadros)
    clear_shape(s20.shapes[9])
    clear_shape(s20.shapes[10])
    print("  ✓ slide 20 reorganizado con cards y limitaciones")

    # ── Slide 21: Resultados y Conclusiones — similar estructura ──
    print("\n═══ Fix slide 21: Resultados y Conclusiones ═══")
    s21 = prs.slides[21 - 1]
    set_shape_lines(s21.shapes[3], ["1"])
    set_shape_lines(s21.shapes[5], ["2"])
    set_shape_lines(s21.shapes[2], ["71 tablas, 245 endpoints"])
    set_shape_lines(s21.shapes[4], ["139 pruebas pasando"])
    set_shape_lines(s21.shapes[1], ["Resultados"])
    set_shape_lines(s21.shapes[6], ["Sistema en producción con HTTPS, observabilidad y CI/CD automatizado desde GitHub vía Dokploy + Traefik."])
    set_shape_lines(s21.shapes[7], ["Es viable acoplar deep learning Python a un backend transaccional Java sin sacrificar estabilidad ni cumplimiento fiscal."])
    set_shape_lines(s21.shapes[8], ["Conclusiones"])
    print("  ✓ slide 21 reorganizado")

    prs.save(OUT)
    print(f"\n✓ Guardado: {OUT}")


if __name__ == "__main__":
    main()
