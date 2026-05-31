"""
Fix-v4 del Sustentacion_Jesus.pptx.

Problemas detectados en v3:
  - Slide 12 (Fase de Diseño): el logo de Figma sigue presente (era image10.png en slide 15 original).
  - Slide 16 (Implementación): shape 5 (mini-label rojo) tiene texto largo
    cortado; debería ser solo el número "1".
  - Slide 17 (Estructura paquetes): la imagen image5 tiene la advertencia
    de PlantUML "Please use CSS style instead of skinparam padding".
    Ya regeneré image5.png sin esa advertencia. Hay que re-insertarla en el slide.
"""
import pathlib
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

ROOT = pathlib.Path("/Users/jesus/Desktop/parqueaderos-api")
OUT = ROOT / "parqueaderos-docs-referencia/SUBIR_A_DRIVE/Sustentacion_Jesus.pptx"
IMG_PAQUETES = ROOT / "docs/figures/main/image5.png"


def set_paragraph_text(paragraph, new_text):
    if not paragraph.runs:
        run_xml = '<a:r xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:rPr lang="es-CO" dirty="0"/><a:t></a:t></a:r>'
        new_r = etree.fromstring(run_xml)
        paragraph._p.append(new_r)
    paragraph.runs[0].text = new_text
    for r in paragraph.runs[1:]:
        r._r.getparent().remove(r._r)


def set_shape_lines(shape, lines):
    if not shape.has_text_frame:
        return
    paras = list(shape.text_frame.paragraphs)
    for i, line in enumerate(lines):
        if i < len(paras):
            set_paragraph_text(paras[i], line)
    for i in range(len(lines), len(paras)):
        set_paragraph_text(paras[i], "")


def remove_picture_in_range(slide, min_w, max_w):
    n = 0
    for shape in list(slide.shapes):
        if shape.shape_type != 13:
            continue
        if min_w <= shape.width <= max_w:
            shape._element.getparent().remove(shape._element)
            n += 1
    return n


def replace_image_in_slide(slide, new_image_path, min_w=None):
    """Encuentra la imagen más grande (filtrando por min_w) y la reemplaza."""
    pictures = [s for s in slide.shapes if s.shape_type == 13]
    if min_w:
        pictures = [p for p in pictures if p.width >= min_w]
    if not pictures:
        return False
    pictures.sort(key=lambda s: s.width * s.height, reverse=True)
    target = pictures[0]
    left, top, width, height = target.left, target.top, target.width, target.height
    sp = target._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(str(new_image_path), left, top, width=width, height=height)
    return True


def main():
    prs = Presentation(OUT)

    # ── Slide 12: eliminar logo Figma (~43cm de ancho) ──
    print("═══ Fix slide 12: Eliminar logo Figma ═══")
    s12 = prs.slides[12 - 1]
    # El logo está en L39 T38 W43 H21 (en unidades /100000 EMU)
    # W=43*100000 = 4_300_000 EMU. Filtramos por ancho entre 3-5 millones EMU.
    n = remove_picture_in_range(s12, min_w=3_000_000, max_w=5_500_000)
    print(f"  ✓ {n} imagen(es) eliminada(s)")

    # ── Slide 16: arreglar mini-label shape 5 ──
    print("\n═══ Fix slide 16: mini-label shape 5 ═══")
    s16 = prs.slides[16 - 1]
    # shape 5 era el mini-label numérico "1" que yo había puesto con texto largo
    set_shape_lines(s16.shapes[5], ["1"])
    print("  ✓ shape 5 → '1'")

    # ── Slide 17: reinsertar image5 (sin warning de PlantUML) ──
    print("\n═══ Fix slide 17: reinsertar image5 sin warning ═══")
    s17 = prs.slides[17 - 1]
    ok = replace_image_in_slide(s17, IMG_PAQUETES, min_w=5_000_000)
    print(f"  ✓ image5 reinsertada: {ok}")

    prs.save(OUT)
    print(f"\n✓ Guardado: {OUT}")


if __name__ == "__main__":
    main()
