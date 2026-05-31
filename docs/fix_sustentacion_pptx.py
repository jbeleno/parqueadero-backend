"""
Correcciones al Sustentacion_Jesus.pptx:
  1. Slide 17 (Estructura de paquetes): vaciar textos viejos de Linda + agregar imagen image5 grande
  2. Slides 22 (Aportes Académicos), 23 (Aportes Profesionales), 24 (Bibliografía):
     reescribir el shape 1 con MULTIPLES párrafos (los aportes / bibliografía van como bullets dentro de un solo text_frame)
  3. Slide 25 (Demo Swagger): si el reemplazo de imagen quedó bien, dejarlo; si no, agregar la imagen
"""
import pathlib
from copy import deepcopy
from pptx import Presentation
from pptx.util import Emu, Cm
from pptx.oxml.ns import qn
from lxml import etree

ROOT = pathlib.Path("/Users/jesus/Desktop/parqueaderos-api")
OUT = ROOT / "parqueaderos-docs-referencia/SUBIR_A_DRIVE/Sustentacion_Jesus.pptx"
IMG_PAQUETES = ROOT / "docs/figures/main/image5.png"


def replace_text_frame(shape, lines):
    """Reemplaza completamente los párrafos del shape con `lines`,
    clonando el formato del primer párrafo para los nuevos."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    paras = list(tf.paragraphs)
    if not paras:
        return
    template_p = paras[0]._p  # plantilla de formato (el primer párrafo)
    # Capturar primer run para clonar estilo
    first_run = template_p.find(qn("a:r"))
    rPr_xml = etree.tostring(first_run.find(qn("a:rPr"))) if first_run is not None and first_run.find(qn("a:rPr")) is not None else b'<a:rPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" lang="es-CO" dirty="0"/>'
    pPr_el = template_p.find(qn("a:pPr"))
    pPr_xml = etree.tostring(pPr_el) if pPr_el is not None else None

    # Remover TODOS los párrafos existentes
    txBody = template_p.getparent()
    for p in list(txBody.findall(qn("a:p"))):
        txBody.remove(p)

    # Agregar uno nuevo por cada línea
    for line in lines:
        new_p = etree.SubElement(txBody, qn("a:p"))
        if pPr_xml is not None:
            new_p.insert(0, etree.fromstring(pPr_xml))
        new_r = etree.SubElement(new_p, qn("a:r"))
        new_r.append(etree.fromstring(rPr_xml))
        new_t = etree.SubElement(new_r, qn("a:t"))
        new_t.text = line


def clear_shape_text(shape):
    """Vacía un text_frame (deja un solo párrafo con texto vacío)."""
    if not shape.has_text_frame:
        return
    paras = list(shape.text_frame.paragraphs)
    if not paras:
        return
    first_p = paras[0]._p
    txBody = first_p.getparent()
    for p in list(txBody.findall(qn("a:p"))):
        txBody.remove(p)
    new_p = etree.SubElement(txBody, qn("a:p"))


# ─────────────────────────────────────────────────────────────────────────────
# Contenido correcto: TODO el bloque va en un solo shape como múltiples párrafos
# ─────────────────────────────────────────────────────────────────────────────

APORTES_ACADEMICOS = [
    "Aplicación integral del SDLC sobre un sistema real, con cuatro anexos formales (Requerimientos IEEE 29148, Historias de Usuario, 42 Diagramas de Casos de Uso UML y 46 ERDs/diagramas de estado/secuencia).",
    "Combinación práctica de tres paradigmas de persistencia coexistentes: relacional clásico, espacial con PostGIS y event sourcing (movimiento_saldo, audit_log), todos sobre PostgreSQL 16.",
    "Demostración empírica del valor de los snapshots de historicidad: 13 columnas *_snapshot en ticket, 4 en factura y 1 en pago preservan la verdad histórica frente a cambios posteriores en entidades referenciadas.",
    "Benchmark reproducible de siete soluciones de OCR para placas colombianas (12 placas × 5 variaciones), con metodología explícita, criterios de evaluación y resultados contraintuitivos sobre los modelos propios entrenados.",
]

APORTES_PROFESIONALES = [
    "Dominio del stack Java empresarial moderno: Spring Boot 3.5, Java 21, Spring Security con JWT y RBAC, Spring Data JPA con Hibernate, Spring WebSocket sobre STOMP, Spring AOP para auditoría y Spring Actuator para observabilidad.",
    "Integración práctica de visión por computador con Python: PyTorch + Ultralytics YOLO para entrenamiento, PaddleOCR para inferencia, FastAPI para exponer el modelo como microservicio y un cliente Java tolerante a fallos.",
    "Manejo completo del ciclo de despliegue moderno: contenedores Docker multi-stage, orquestación con docker-compose, infraestructura Dokploy con Traefik para HTTPS, CI/CD con GitHub Actions y registro de imágenes en GHCR.",
    "Práctica de documentación técnica formal: validación del schema contra producción con psql, glosarios de nomenclatura, fichas de RF con plantilla IEEE 29148 y trazabilidad completa entre RF-HU-UC-pruebas.",
]

BIBLIOGRAFIA = [
    "Spring Team. (2024). Spring Boot Reference Documentation v3.5. https://docs.spring.io/spring-boot/",
    "The PostgreSQL Global Development Group. (2024). PostgreSQL 16 Documentation. https://www.postgresql.org/docs/16/",
    "Ultralytics. (2024). YOLOv11 Documentation. https://docs.ultralytics.com/",
    "PaddlePaddle. (2024). PaddleOCR Documentation. https://paddlepaddle.github.io/PaddleOCR/",
    "IEEE/ISO/IEC. (2018). 29148-2018 — Systems and software engineering — Life cycle processes — Requirements engineering.",
    "DIAN. (2023). Resoluciones de facturación electrónica. https://www.dian.gov.co/",
    "Fowler, M. (2003). Patterns of Enterprise Application Architecture. Addison-Wesley.",
    "Gamma, E., Helm, R., Johnson, R., Vlissides, J. (1995). Design Patterns. Addison-Wesley.",
]


def main():
    prs = Presentation(OUT)

    # ── Slide 17: vaciar textos viejos y agregar imagen ──
    print("═══ Fix slide 17: Estructura de paquetes ═══")
    s17 = prs.slides[17 - 1]
    for j, shape in enumerate(s17.shapes):
        if j == 0:
            continue  # mantenemos título
        if shape.has_text_frame:
            clear_shape_text(shape)
            print(f"  ✓ vaciado shape {j}")
    # Agregar imagen centrada, ocupando el 70% del slide (debajo del título)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    # Imagen image5: 2898x1129 px → ratio 2.57
    img_ratio = 2898 / 1129
    img_w = int(slide_w * 0.78)
    img_h = int(img_w / img_ratio)
    if img_h > slide_h * 0.7:
        img_h = int(slide_h * 0.7)
        img_w = int(img_h * img_ratio)
    left = (slide_w - img_w) // 2
    top = int(slide_h * 0.25)
    s17.shapes.add_picture(str(IMG_PAQUETES), left, top, width=img_w, height=img_h)
    print(f"  ✓ image5.png agregada ({img_w}x{img_h} EMU)")

    # ── Slide 22: Aportes Académicos en shape 1 ──
    print("\n═══ Fix slide 22: Aportes Académicos ═══")
    s22 = prs.slides[22 - 1]
    replace_text_frame(s22.shapes[1], APORTES_ACADEMICOS)
    print(f"  ✓ {len(APORTES_ACADEMICOS)} aportes académicos")

    # ── Slide 23: Aportes Profesionales en shape 1 ──
    print("\n═══ Fix slide 23: Aportes Profesionales ═══")
    s23 = prs.slides[23 - 1]
    replace_text_frame(s23.shapes[1], APORTES_PROFESIONALES)
    print(f"  ✓ {len(APORTES_PROFESIONALES)} aportes profesionales")

    # ── Slide 24: Bibliografía en shape 1 ──
    print("\n═══ Fix slide 24: Bibliografía ═══")
    s24 = prs.slides[24 - 1]
    replace_text_frame(s24.shapes[1], BIBLIOGRAFIA)
    print(f"  ✓ {len(BIBLIOGRAFIA)} referencias")

    prs.save(OUT)
    print(f"\n✓ Guardado: {OUT}")


if __name__ == "__main__":
    main()
