"""
Fix-v8:
  1. Recuperar formato de título (azul USCO #003462, sz=3000, b=1) en
     slides 11, 13, 14, 15 — se perdió al setear texto en runs sin rPr.
  2. Re-insertar image3 (ERD nuevo con ratio 0.80) en slide 13 con tamaño
     adecuado.
"""
import pathlib
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image

ROOT = pathlib.Path("/Users/jesus/Desktop/parqueaderos-api")
OUT = ROOT / "parqueaderos-docs-referencia/SUBIR_A_DRIVE/presentacion/pptx/Sustentacion_Jesus.pptx"
FIG = ROOT / "docs/figures/main"


TITLE_RPR_XML = (
    '<a:rPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
    'lang="es-ES" sz="3000" b="1">'
    '<a:solidFill><a:srgbClr val="003462"/></a:solidFill>'
    '</a:rPr>'
)


def apply_title_format(shape):
    """Aplica formato azul USCO sz=3000 b=1 al primer párrafo del shape."""
    if not shape.has_text_frame:
        return
    paras = list(shape.text_frame.paragraphs)
    if not paras:
        return
    p = paras[0]._p
    # Buscar el primer run, agregar/reemplazar rPr
    run = p.find(qn("a:r"))
    if run is None:
        return
    # Eliminar rPr existente si hay
    existing_rPr = run.find(qn("a:rPr"))
    if existing_rPr is not None:
        run.remove(existing_rPr)
    # Insertar el rPr nuevo al inicio del run
    new_rPr = etree.fromstring(TITLE_RPR_XML)
    run.insert(0, new_rPr)


def replace_picture_centered(slide, img_path, slide_w, slide_h):
    """Elimina las imágenes anteriores y agrega img_path centrada en el área
    disponible (debajo del banner/título, hasta el footer)."""
    for shape in list(slide.shapes):
        if shape.shape_type == 13:
            shape._element.getparent().remove(shape._element)
    with Image.open(img_path) as im:
        ratio = im.width / im.height
    top_avail = int(slide_h * 0.18)
    bottom_avail = int(slide_h * 0.96)
    height_avail = bottom_avail - top_avail
    max_width = int(slide_w * 0.95)
    new_h = height_avail
    new_w = int(new_h * ratio)
    if new_w > max_width:
        new_w = max_width
        new_h = int(new_w / ratio)
    left = (slide_w - new_w) // 2
    top = top_avail + (height_avail - new_h) // 2
    slide.shapes.add_picture(str(img_path), left, top, width=new_w, height=new_h)
    return (left, top, new_w, new_h)


def main():
    prs = Presentation(OUT)
    sw, sh = prs.slide_width, prs.slide_height

    # ── Aplicar formato del título en slides 11, 13, 14, 15 ──
    print("═══ Formato azul USCO al título de slides 11, 13, 14, 15 ═══")
    for idx in [11, 13, 14, 15]:
        slide = prs.slides[idx - 1]
        apply_title_format(slide.shapes[0])
        # Obtener texto del título para verificar
        title_txt = slide.shapes[0].text_frame.paragraphs[0].text
        print(f"  slide {idx}: '{title_txt}' → azul #003462, 30pt, bold")

    # ── Re-insertar image3 ERD con su nueva versión más cuadrada ──
    print("\n═══ Re-insertar image3 ERD en slide 13 (ratio 0.80) ═══")
    s13 = prs.slides[13 - 1]
    pos = replace_picture_centered(s13, FIG / "image3.png", sw, sh)
    print(f"  → L{pos[0]//100000} T{pos[1]//100000} W{pos[2]//100000} H{pos[3]//100000}")

    prs.save(OUT)
    print(f"\n✓ Guardado: {OUT}")


if __name__ == "__main__":
    main()
