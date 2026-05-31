"""
Fix-v9: dividir el diagrama de casos de uso del slide 11 en DOS imágenes
lado a lado:
  - image1a.png (principal): Autenticación + Operación + Reservas + Flujo automático
  - image1b.png (administración): ADMIN + SUPER_ADMIN + 4 casos de uso

Layout deseado en el slide 11:
  ┌──────────────────────────────────┐
  │  Título azul USCO arriba         │
  │                                  │
  │  ┌─────────────┐  ┌───────────┐  │
  │  │             │  │           │  │
  │  │  image1a    │  │ image1b   │  │
  │  │ (principal) │  │  (admin)  │  │
  │  │             │  │           │  │
  │  └─────────────┘  └───────────┘  │
  └──────────────────────────────────┘
"""
import pathlib
from pptx import Presentation
from PIL import Image

ROOT = pathlib.Path("/Users/jesus/Desktop/parqueaderos-api")
OUT = ROOT / "parqueaderos-docs-referencia/SUBIR_A_DRIVE/presentacion/pptx/Sustentacion_Jesus.pptx"
FIG = ROOT / "docs/figures/main"
IMG_A = FIG / "image1a.png"  # principal
IMG_B = FIG / "image1b.png"  # admin


def main():
    prs = Presentation(OUT)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    s11 = prs.slides[11 - 1]

    # Eliminar todas las imágenes actuales del slide
    for shape in list(s11.shapes):
        if shape.shape_type == 13:
            shape._element.getparent().remove(shape._element)

    # Área disponible para las dos imágenes (debajo de título+banner, encima de footer)
    top_avail = int(slide_h * 0.18)
    bottom_avail = int(slide_h * 0.96)
    height_avail = bottom_avail - top_avail

    # División horizontal: dejamos un gap entre las dos imágenes
    gap = int(slide_w * 0.02)
    side_margin = int(slide_w * 0.03)
    total_w = slide_w - 2 * side_margin - gap
    # Damos 65% del ancho a la principal y 35% a la de admin
    w_a_max = int(total_w * 0.62)
    w_b_max = total_w - w_a_max

    # Calcular tamaños respetando aspect ratio
    with Image.open(IMG_A) as im_a:
        ratio_a = im_a.width / im_a.height
    with Image.open(IMG_B) as im_b:
        ratio_b = im_b.width / im_b.height

    # Probar limitado por altura primero
    h_a = height_avail
    w_a = int(h_a * ratio_a)
    if w_a > w_a_max:
        w_a = w_a_max
        h_a = int(w_a / ratio_a)

    h_b = height_avail
    w_b = int(h_b * ratio_b)
    if w_b > w_b_max:
        w_b = w_b_max
        h_b = int(w_b / ratio_b)

    # Centrado vertical de cada imagen en el área disponible
    top_a = top_avail + (height_avail - h_a) // 2
    top_b = top_avail + (height_avail - h_b) // 2

    # Posición horizontal: image1a a la izquierda, image1b a la derecha
    # Centramos el conjunto de las dos
    total_used_w = w_a + gap + w_b
    start_left = (slide_w - total_used_w) // 2
    left_a = start_left
    left_b = start_left + w_a + gap

    s11.shapes.add_picture(str(IMG_A), left_a, top_a, width=w_a, height=h_a)
    s11.shapes.add_picture(str(IMG_B), left_b, top_b, width=w_b, height=h_b)

    print(f"image1a: L{left_a//100000} T{top_a//100000} W{w_a//100000} H{h_a//100000}")
    print(f"image1b: L{left_b//100000} T{top_b//100000} W{w_b//100000} H{h_b//100000}")

    prs.save(OUT)
    print(f"\n✓ Guardado: {OUT}")


if __name__ == "__main__":
    main()
