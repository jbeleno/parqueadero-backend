"""
Fix-v5: arreglar layout de slides 20 y 21 según estructura real de Linda.

Layout descubierto al analizar coordenadas:

  Slide 20 (Alcance y Limitaciones):
    shape 0: título grande arriba
    shape 3 ("1") + shape 2 (texto corto al lado) → mini-stat izquierda arriba
    shape 5 ("2") + shape 4 (texto corto al lado) → mini-stat derecha arriba
    shape 1 (T=23, W=56) → header "Alcance"
    shape 6 (T=34, W=56) → bullet alcance 1
    shape 7 (T=44, W=56, H=17) → bullet alcance 2  (H grande)
    shape 8 (T=23, W=56) → header "Limitaciones"
    shape 9 (T=34, W=56) → bullet limit 1
    shape 10 (T=44, W=56) → bullet limit 2

  Slide 21 (Resultados y Conclusiones):
    Misma estructura pero con SOLO 9 shapes en lugar de 11.
    shape 8 (T=23, W=56, H=17) reemplaza a shapes 8+9+10 de slide 20:
    debe contener "Conclusiones" + bullets como multi-paragraph.
"""
import pathlib
from copy import deepcopy
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

ROOT = pathlib.Path("/Users/jesus/Desktop/parqueaderos-api")
OUT = ROOT / "parqueaderos-docs-referencia/SUBIR_A_DRIVE/Sustentacion_Jesus.pptx"


def set_paragraph_text(paragraph, new_text):
    if not paragraph.runs:
        run_xml = '<a:r xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:rPr lang="es-CO" dirty="0"/><a:t></a:t></a:r>'
        new_r = etree.fromstring(run_xml)
        paragraph._p.append(new_r)
    paragraph.runs[0].text = new_text
    for r in paragraph.runs[1:]:
        r._r.getparent().remove(r._r)


def set_shape_lines(shape, lines):
    if not shape.has_text_frame:
        return
    paras = list(shape.text_frame.paragraphs)
    for i, line in enumerate(lines):
        if i < len(paras):
            set_paragraph_text(paras[i], line)
    for i in range(len(lines), len(paras)):
        set_paragraph_text(paras[i], "")


def replace_text_frame_multi(shape, lines):
    """Reemplaza completamente el text_frame con N párrafos clonando formato del primero."""
    if not shape.has_text_frame:
        return
    paras = list(shape.text_frame.paragraphs)
    if not paras:
        return
    template_p = paras[0]._p
    first_run = template_p.find(qn("a:r"))
    if first_run is not None and first_run.find(qn("a:rPr")) is not None:
        rPr_xml = etree.tostring(first_run.find(qn("a:rPr")))
    else:
        rPr_xml = b'<a:rPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" lang="es-CO" dirty="0"/>'
    pPr_el = template_p.find(qn("a:pPr"))
    pPr_xml = etree.tostring(pPr_el) if pPr_el is not None else None
    txBody = template_p.getparent()
    for p in list(txBody.findall(qn("a:p"))):
        txBody.remove(p)
    for line in lines:
        new_p = etree.SubElement(txBody, qn("a:p"))
        if pPr_xml is not None:
            new_p.insert(0, etree.fromstring(pPr_xml))
        new_r = etree.SubElement(new_p, qn("a:r"))
        new_r.append(etree.fromstring(rPr_xml))
        new_t = etree.SubElement(new_r, qn("a:t"))
        new_t.text = line


def main():
    prs = Presentation(OUT)

    # ── Slide 20: Alcance y Limitaciones ──
    print("═══ Fix slide 20: Alcance y Limitaciones ═══")
    s20 = prs.slides[20 - 1]
    set_shape_lines(s20.shapes[1], ["Alcance"])
    set_shape_lines(s20.shapes[2], ["245 endpoints"])
    set_shape_lines(s20.shapes[3], ["1"])
    set_shape_lines(s20.shapes[4], ["Sidecar OCR funcional"])
    set_shape_lines(s20.shapes[5], ["2"])
    set_shape_lines(s20.shapes[6], ["Backend con multi-tenant, tickets, reservas, facturación DIAN, suscripciones, convenios y auditoría."])
    set_shape_lines(s20.shapes[7], ["YOLOv11n con mAP50=0.987 y PaddleOCR v2 con 100% de consistencia sobre 60 imágenes de prueba."])
    set_shape_lines(s20.shapes[8], ["Limitaciones"])
    set_shape_lines(s20.shapes[9], ["El alcance es backend + sidecar OCR; el cliente móvil y web son desarrollos paralelos no incluidos."])
    set_shape_lines(s20.shapes[10], ["Montos en double precision (anti-patrón); campos DIAN nuevos en numeric(14,2). Migración pendiente."])
    print("  ✓ slide 20 reorganizado con bullets alcance + limitaciones")

    # ── Slide 21: Resultados y Conclusiones ──
    print("\n═══ Fix slide 21: Resultados y Conclusiones ═══")
    s21 = prs.slides[21 - 1]
    set_shape_lines(s21.shapes[1], ["Resultados"])
    set_shape_lines(s21.shapes[2], ["71 tablas, 245 endpoints"])
    set_shape_lines(s21.shapes[3], ["1"])
    set_shape_lines(s21.shapes[4], ["139 pruebas pasando"])
    set_shape_lines(s21.shapes[5], ["2"])
    set_shape_lines(s21.shapes[6], ["Sistema en producción con HTTPS, observabilidad y CI/CD automatizado desde GitHub vía Dokploy + Traefik."])
    set_shape_lines(s21.shapes[7], ["YOLOv11n entrenado a propósito (mAP50=0.987) con consistencia del 100% sobre el benchmark de 60 imágenes."])
    # shape 8 (H=17, alto) → multi-paragraph: header "Conclusiones" + 3 bullets
    replace_text_frame_multi(s21.shapes[8], [
        "Conclusiones",
        "Es viable acoplar deep learning Python a un backend transaccional Java sin sacrificar estabilidad.",
        "El Modelo B + IVA + tope + resoluciones DIAN dejan al sistema listo para certificación fiscal.",
        "Auditoría AOP y RBAC granular soportan trazabilidad y aislamiento multi-tenant.",
    ])
    print("  ✓ slide 21 reorganizado con bullets resultados + conclusiones (multi-paragraph)")

    prs.save(OUT)
    print(f"\n✓ Guardado: {OUT}")


if __name__ == "__main__":
    main()
