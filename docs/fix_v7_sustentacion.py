"""
Fix-v7: reposicionar imágenes de slides 11, 13, 14, 15.

Problema: las 4 imágenes están en T=3 H=64, lo que las hace ocupar casi
todo el slide y TAPA el título (que está en L=66 T=3) y también el banner
decorativo image3.

Solución: bajar la imagen a T=14 (debajo del banner+título) y dimensionarla
manteniendo el ratio de la imagen original, centrada horizontalmente.

Slide: 121cm × 68cm
  - Banner superior decorativo: T=0 a ~T=10
  - Título (texto): L=66 T=3 W=53 H=8 → ocupa esquina superior derecha
  - Banner inferior: ~T=66 a T=68
  - Área libre para imagen: T=13 a T=65 → altura disponible ≈ 52cm
"""
import pathlib
from pptx import Presentation
from PIL import Image

ROOT = pathlib.Path("/Users/jesus/Desktop/parqueaderos-api")
OUT = ROOT / "parqueaderos-docs-referencia/SUBIR_A_DRIVE/presentacion/pptx/Sustentacion_Jesus.pptx"
FIG = ROOT / "docs/figures/main"

# Mapeo slide → ruta de la imagen actual del slide
IMG_BY_SLIDE = {
    11: FIG / "image1.png",   # Casos de uso
    13: FIG / "image3.png",   # ERD jerárquico
    14: FIG / "image2.png",   # Arquitectura
    15: FIG / "image8.png",   # Estados Ticket
}


def main():
    prs = Presentation(OUT)
    slide_w = prs.slide_width   # 12192000 EMU
    slide_h = prs.slide_height  # 6858000 EMU

    # Área disponible para la imagen (en EMU)
    top_avail   = int(slide_h * 0.18)   # T ≈ 12cm (debajo de título y banner)
    bottom_avail = int(slide_h * 0.96)  # B ≈ 65cm (encima del footer)
    height_avail = bottom_avail - top_avail

    max_width = int(slide_w * 0.95)  # 95% del ancho del slide

    for slide_idx, img_path in IMG_BY_SLIDE.items():
        slide = prs.slides[slide_idx - 1]
        # Eliminar imagen actual (la que tapa todo)
        for shape in list(slide.shapes):
            if shape.shape_type == 13:
                shape._element.getparent().remove(shape._element)

        # Calcular tamaño manteniendo aspect ratio
        with Image.open(img_path) as im:
            ratio = im.width / im.height
        # Empezar limitado por altura
        new_h = height_avail
        new_w = int(new_h * ratio)
        # Si excede el ancho, ajustar por ancho
        if new_w > max_width:
            new_w = max_width
            new_h = int(new_w / ratio)
        # Centrar
        left = (slide_w - new_w) // 2
        top = top_avail + (height_avail - new_h) // 2  # centrar verticalmente en el área disponible

        slide.shapes.add_picture(str(img_path), left, top, width=new_w, height=new_h)
        print(f"  slide {slide_idx}: {img_path.name} → "
              f"L{left//100000} T{top//100000} W{new_w//100000} H{new_h//100000} "
              f"(ratio {ratio:.2f})")

    prs.save(OUT)
    print(f"\n✓ Guardado: {OUT}")


if __name__ == "__main__":
    main()
