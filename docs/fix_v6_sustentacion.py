"""
Fix-v6 final: arreglar slide 25 (Demo Swagger).

Problema: la imagen image4.png está en posición L=48 T=11 W=25 H=55 (en cm aprox)
es decir, está pegada a la mitad-derecha y muy chica (25cm de ancho).

Solución: agrandar la imagen y centrarla horizontalmente debajo del título.
"""
import pathlib
from pptx import Presentation
from PIL import Image

ROOT = pathlib.Path("/Users/jesus/Desktop/parqueaderos-api")
OUT = ROOT / "parqueaderos-docs-referencia/SUBIR_A_DRIVE/Sustentacion_Jesus.pptx"
IMG_SWAGGER = ROOT / "docs/figures/main/image4.png"


def main():
    prs = Presentation(OUT)
    s25 = prs.slides[25 - 1]
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Eliminar imagen actual
    for shape in list(s25.shapes):
        if shape.shape_type == 13:  # picture
            shape._element.getparent().remove(shape._element)
            print("  ✓ imagen vieja eliminada")

    # Calcular dimensiones para image4 centrada (ocupa 75% del ancho, manteniendo ratio)
    with Image.open(IMG_SWAGGER) as im:
        ratio = im.width / im.height
    img_w = int(slide_w * 0.55)
    img_h = int(img_w / ratio)
    # Si la altura supera el 70% del slide, ajustamos por altura
    if img_h > slide_h * 0.72:
        img_h = int(slide_h * 0.72)
        img_w = int(img_h * ratio)
    left = (slide_w - img_w) // 2
    top = int(slide_h * 0.20)  # debajo del título

    s25.shapes.add_picture(str(IMG_SWAGGER), left, top, width=img_w, height=img_h)
    print(f"  ✓ image4 reinsertada centrada ({img_w//100000}×{img_h//100000} cm)")

    prs.save(OUT)
    print(f"\n✓ Guardado: {OUT}")


if __name__ == "__main__":
    main()
